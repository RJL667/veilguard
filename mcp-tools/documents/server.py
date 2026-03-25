"""Document MCP Tool Server for Veilguard.

Tools: read_pdf, create_pdf, read_docx, create_docx, read_xlsx, create_xlsx, read_pptx, create_pptx
"""

import json
import os
from pathlib import Path

from mcp.server.fastmcp import FastMCP

mcp = FastMCP(
    "documents",
    instructions="Document tools for reading and creating PDF, Word, Excel, and PowerPoint files.",
)

WORKSPACE = os.environ.get("WORKSPACE_ROOT", "/workspace")


def _safe_path(path: str) -> Path:
    """Resolve path relative to workspace, prevent escape."""
    p = Path(path)
    if not p.is_absolute():
        p = Path(WORKSPACE) / p
    resolved = p.resolve()
    workspace_resolved = Path(WORKSPACE).resolve()
    if not str(resolved).startswith(str(workspace_resolved)):
        raise ValueError(f"Access denied: path must be within {WORKSPACE}")
    return resolved


# ── PDF ──────────────────────────────────────────────────────────────────────


@mcp.tool()
def read_pdf(path: str, pages: str = "") -> str:
    """Extract text from a PDF file.

    Args:
        path: PDF file path (absolute or relative to workspace)
        pages: Page range to extract, e.g. "1-5" or "3". Empty = all pages.
    """
    import fitz  # PyMuPDF

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        doc = fitz.open(str(p))
        total = len(doc)

        # Parse page range
        if pages.strip():
            parts = pages.strip().split("-")
            start = max(0, int(parts[0]) - 1)
            end = int(parts[-1]) if len(parts) > 1 else start + 1
            end = min(end, total)
            page_range = range(start, end)
        else:
            page_range = range(total)

        output = [f"# {p.name} ({total} pages, reading {len(page_range)})\n"]
        for i in page_range:
            page = doc[i]
            text = page.get_text().strip()
            output.append(f"--- Page {i + 1} ---\n{text}\n")

        doc.close()
        return "\n".join(output)
    except Exception as e:
        return f"Error reading PDF: {e}"


@mcp.tool()
def create_pdf(path: str, content: str, title: str = "") -> str:
    """Create a PDF file from text content.

    Args:
        path: Output PDF path (absolute or relative to workspace)
        content: Text content for the PDF. Supports basic formatting with lines.
        title: Optional title for the first page
    """
    import fitz  # PyMuPDF

    p = _safe_path(path)
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = fitz.open()

        # Simple text-to-PDF: split content into pages (~60 lines each)
        lines = content.splitlines()
        if title:
            lines.insert(0, "")
            lines.insert(0, title)
            lines.insert(1, "=" * len(title))
            lines.insert(2, "")

        lines_per_page = 55
        font_size = 11
        margin = 72  # 1 inch

        for page_start in range(0, len(lines), lines_per_page):
            page = doc.new_page()
            page_lines = lines[page_start : page_start + lines_per_page]
            text_block = "\n".join(page_lines)

            # Insert text
            rect = fitz.Rect(margin, margin, page.rect.width - margin, page.rect.height - margin)
            page.insert_textbox(
                rect,
                text_block,
                fontsize=font_size,
                fontname="helv",
            )

        doc.save(str(p))
        doc.close()
        return f"Created PDF: {p} ({len(doc) if doc else (len(lines) // lines_per_page + 1)} pages)"
    except Exception as e:
        return f"Error creating PDF: {e}"


# ── PDF Edit ─────────────────────────────────────────────────────────────────


@mcp.tool()
def edit_pdf_text(path: str, old_text: str, new_text: str, page: int = 0) -> str:
    """Find and replace text in a PDF. Searches all pages or a specific page.

    Note: PDF text replacement works best with simple text. Complex layouts may not
    replace cleanly — in that case, use redact + insert approach.

    Args:
        path: PDF file path
        old_text: Text to find
        new_text: Replacement text
        page: Page number (1-indexed). 0 = search all pages.
    """
    import fitz

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        doc = fitz.open(str(p))
        count = 0
        pages_to_search = [page - 1] if page > 0 else range(len(doc))

        for pg_idx in pages_to_search:
            if pg_idx >= len(doc):
                continue
            pg = doc[pg_idx]
            # Find all instances
            instances = pg.search_for(old_text)
            for inst in instances:
                # Redact old text
                pg.add_redact_annot(inst, new_text, fontsize=11, fontname="helv")
                count += 1
            if instances:
                pg.apply_redactions()

        if count == 0:
            doc.close()
            return f"Error: '{old_text}' not found in {p.name}"

        doc.save(str(p), incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP)
        doc.close()
        return f"Replaced {count} occurrence(s) of '{old_text}' in {p.name}"
    except Exception as e:
        return f"Error editing PDF: {e}"


@mcp.tool()
def pdf_add_text(path: str, text: str, page: int = 1, x: float = 72, y: float = 72, fontsize: float = 12) -> str:
    """Add text to a specific position on a PDF page.

    Args:
        path: PDF file path
        text: Text to add
        page: Page number (1-indexed)
        x: X position in points from left edge (72 = 1 inch)
        y: Y position in points from top edge (72 = 1 inch)
        fontsize: Font size (default 12)
    """
    import fitz

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        doc = fitz.open(str(p))
        if page < 1 or page > len(doc):
            doc.close()
            return f"Error: Page {page} out of range (1-{len(doc)})"

        pg = doc[page - 1]
        point = fitz.Point(x, y)
        pg.insert_text(point, text, fontsize=fontsize, fontname="helv")

        doc.save(str(p), incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP)
        doc.close()
        return f"Added text to page {page} at ({x}, {y}) in {p.name}"
    except Exception as e:
        return f"Error adding text: {e}"


@mcp.tool()
def pdf_delete_pages(path: str, pages: str) -> str:
    """Delete pages from a PDF.

    Args:
        path: PDF file path
        pages: Pages to delete. Comma-separated or range: "3", "1,3,5", "2-4"
    """
    import fitz

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        doc = fitz.open(str(p))
        total = len(doc)

        # Parse page numbers
        to_delete = set()
        for part in pages.split(","):
            part = part.strip()
            if "-" in part:
                start, end = part.split("-", 1)
                to_delete.update(range(int(start), int(end) + 1))
            else:
                to_delete.add(int(part))

        # Convert to 0-indexed and sort descending (delete from end)
        indices = sorted([i - 1 for i in to_delete if 1 <= i <= total], reverse=True)
        if not indices:
            doc.close()
            return f"Error: No valid pages to delete (file has {total} pages)"

        for idx in indices:
            doc.delete_page(idx)

        doc.save(str(p), garbage=4, deflate=True)
        doc.close()
        return f"Deleted {len(indices)} page(s) from {p.name} ({total} → {total - len(indices)} pages)"
    except Exception as e:
        return f"Error deleting pages: {e}"


@mcp.tool()
def pdf_merge(paths: str, output_path: str) -> str:
    """Merge multiple PDFs into one.

    Args:
        paths: Comma-separated PDF paths to merge (in order)
        output_path: Output PDF path
    """
    import fitz

    out = _safe_path(output_path)

    try:
        result = fitz.open()
        file_list = [p.strip() for p in paths.split(",")]
        total_pages = 0

        for fp in file_list:
            src = _safe_path(fp)
            if not src.exists():
                result.close()
                return f"Error: File not found: {src}"
            src_doc = fitz.open(str(src))
            result.insert_pdf(src_doc)
            total_pages += len(src_doc)
            src_doc.close()

        out.parent.mkdir(parents=True, exist_ok=True)
        result.save(str(out), garbage=4, deflate=True)
        result.close()
        return f"Merged {len(file_list)} PDFs → {out} ({total_pages} pages)"
    except Exception as e:
        return f"Error merging PDFs: {e}"


@mcp.tool()
def pdf_add_watermark(path: str, text: str = "CONFIDENTIAL", opacity: float = 0.3) -> str:
    """Add a diagonal watermark to every page of a PDF.

    Args:
        path: PDF file path
        text: Watermark text (default "CONFIDENTIAL")
        opacity: Opacity 0.0-1.0 (default 0.3)
    """
    import fitz

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        doc = fitz.open(str(p))
        for pg in doc:
            rect = pg.rect
            # Diagonal across page
            center = fitz.Point(rect.width / 2, rect.height / 2)
            pg.insert_text(
                center,
                text,
                fontsize=60,
                fontname="helv",
                color=(0.8, 0.8, 0.8),
                rotate=45,
                overlay=True,
            )

        doc.save(str(p), incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP)
        doc.close()
        return f"Added watermark '{text}' to {len(doc)} pages in {p.name}"
    except Exception as e:
        return f"Error adding watermark: {e}"


# ── Word (DOCX) ─────────────────────────────────────────────────────────────


@mcp.tool()
def read_docx(path: str) -> str:
    """Extract text from a Word document.

    Args:
        path: DOCX file path (absolute or relative to workspace)
    """
    from docx import Document

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        doc = Document(str(p))
        output = [f"# {p.name}\n"]

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            if "Heading 1" in style:
                output.append(f"\n# {text}")
            elif "Heading 2" in style:
                output.append(f"\n## {text}")
            elif "Heading 3" in style:
                output.append(f"\n### {text}")
            else:
                output.append(text)

        # Extract tables
        for i, table in enumerate(doc.tables):
            output.append(f"\n--- Table {i + 1} ---")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                output.append(" | ".join(cells))

        return "\n".join(output)
    except Exception as e:
        return f"Error reading DOCX: {e}"


@mcp.tool()
def create_docx(path: str, content: str, title: str = "") -> str:
    """Create a Word document from text content.

    Args:
        path: Output DOCX path (absolute or relative to workspace)
        content: Text content. Lines starting with # are headings.
        title: Optional document title
    """
    from docx import Document

    p = _safe_path(path)
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        doc = Document()

        if title:
            doc.add_heading(title, level=0)

        for line in content.splitlines():
            stripped = line.strip()
            if not stripped:
                doc.add_paragraph("")
            elif stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:], style="List Bullet")
            elif len(stripped) > 2 and stripped[0].isdigit() and stripped[1] in ".)" :
                doc.add_paragraph(stripped[2:].strip(), style="List Number")
            else:
                doc.add_paragraph(stripped)

        doc.save(str(p))
        return f"Created DOCX: {p}"
    except Exception as e:
        return f"Error creating DOCX: {e}"


# ── Excel (XLSX) ─────────────────────────────────────────────────────────────


@mcp.tool()
def read_xlsx(path: str, sheet: str = "", max_rows: int = 500) -> str:
    """Read an Excel spreadsheet as text.

    Args:
        path: XLSX file path (absolute or relative to workspace)
        sheet: Sheet name to read. Empty = first sheet.
        max_rows: Maximum rows to return (default 500)
    """
    from openpyxl import load_workbook

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        wb = load_workbook(str(p), read_only=True, data_only=True)

        if sheet:
            if sheet not in wb.sheetnames:
                return f"Error: Sheet '{sheet}' not found. Available: {wb.sheetnames}"
            ws = wb[sheet]
        else:
            ws = wb.active

        output = [f"# {p.name} — Sheet: {ws.title}"]
        output.append(f"# Sheets: {wb.sheetnames}\n")

        rows_read = 0
        for row in ws.iter_rows(values_only=True):
            if rows_read >= max_rows:
                output.append(f"\n... truncated at {max_rows} rows")
                break
            cells = [str(c) if c is not None else "" for c in row]
            output.append(" | ".join(cells))
            rows_read += 1

        wb.close()
        return "\n".join(output)
    except Exception as e:
        return f"Error reading XLSX: {e}"


@mcp.tool()
def create_xlsx(path: str, data: str, sheet_name: str = "Sheet1") -> str:
    """Create an Excel spreadsheet from data.

    Args:
        path: Output XLSX path (absolute or relative to workspace)
        data: JSON string — either a list of dicts [{"col": "val", ...}] or a list of lists [["h1","h2"],["v1","v2"]]
        sheet_name: Name for the worksheet
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font

    p = _safe_path(path)
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        parsed = json.loads(data)
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name

        if not parsed:
            wb.save(str(p))
            return f"Created empty XLSX: {p}"

        if isinstance(parsed[0], dict):
            # List of dicts — keys become headers
            headers = list(parsed[0].keys())
            ws.append(headers)
            for cell in ws[1]:
                cell.font = Font(bold=True)
            for row in parsed:
                ws.append([row.get(h, "") for h in headers])
        elif isinstance(parsed[0], (list, tuple)):
            # List of lists — first row is header
            for i, row in enumerate(parsed):
                ws.append(list(row))
                if i == 0:
                    for cell in ws[1]:
                        cell.font = Font(bold=True)

        # Auto-width columns
        for col in ws.columns:
            max_len = max(len(str(cell.value or "")) for cell in col)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)

        wb.save(str(p))
        row_count = ws.max_row - 1  # minus header
        return f"Created XLSX: {p} ({row_count} rows, {ws.max_column} columns)"
    except json.JSONDecodeError:
        return "Error: 'data' must be a valid JSON string (list of dicts or list of lists)"
    except Exception as e:
        return f"Error creating XLSX: {e}"


# ── Excel Edit ────────────────────────────────────────────────────────────────


@mcp.tool()
def edit_xlsx_cell(path: str, cell: str, value: str, sheet: str = "") -> str:
    """Edit a specific cell in an Excel spreadsheet.

    Args:
        path: XLSX file path
        cell: Cell reference (e.g. "A1", "B5", "C12")
        value: New value. Numbers are auto-detected.
        sheet: Sheet name. Empty = active sheet.
    """
    from openpyxl import load_workbook

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        wb = load_workbook(str(p))
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active

        old_val = ws[cell].value
        # Auto-detect number
        try:
            ws[cell] = float(value) if "." in value else int(value)
        except (ValueError, TypeError):
            ws[cell] = value

        wb.save(str(p))
        return f"Updated {ws.title}!{cell}: '{old_val}' → '{value}'"
    except Exception as e:
        return f"Error editing cell: {e}"


@mcp.tool()
def edit_xlsx_range(path: str, start_cell: str, data: str, sheet: str = "") -> str:
    """Write data to a range of cells starting from a given cell.

    Args:
        path: XLSX file path
        start_cell: Top-left cell (e.g. "A1", "B3")
        data: JSON string — list of lists: [["a","b"],["c","d"]]
        sheet: Sheet name. Empty = active sheet.
    """
    from openpyxl import load_workbook
    from openpyxl.utils import column_index_from_string, get_column_letter

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        parsed = json.loads(data)
        wb = load_workbook(str(p))
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active

        # Parse start cell
        col_str = "".join(c for c in start_cell if c.isalpha())
        row_num = int("".join(c for c in start_cell if c.isdigit()))
        col_num = column_index_from_string(col_str)

        cells_written = 0
        for r, row_data in enumerate(parsed):
            for c, val in enumerate(row_data):
                cell = ws.cell(row=row_num + r, column=col_num + c)
                try:
                    cell.value = float(val) if isinstance(val, str) and "." in val else int(val) if isinstance(val, str) and val.isdigit() else val
                except (ValueError, TypeError):
                    cell.value = val
                cells_written += 1

        wb.save(str(p))
        return f"Updated {cells_written} cells starting at {ws.title}!{start_cell}"
    except json.JSONDecodeError:
        return "Error: 'data' must be a valid JSON list of lists"
    except Exception as e:
        return f"Error editing range: {e}"


# ── Word Edit ─────────────────────────────────────────────────────────────────


@mcp.tool()
def edit_docx(path: str, old_text: str, new_text: str, replace_all: bool = False) -> str:
    """Find and replace text in a Word document.

    Args:
        path: DOCX file path
        old_text: Text to find
        new_text: Replacement text
        replace_all: Replace all occurrences (default: first only)
    """
    from docx import Document

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        doc = Document(str(p))
        count = 0

        for para in doc.paragraphs:
            if old_text in para.text:
                for run in para.runs:
                    if old_text in run.text:
                        if replace_all:
                            run.text = run.text.replace(old_text, new_text)
                            count += run.text.count(new_text) or 1
                        else:
                            run.text = run.text.replace(old_text, new_text, 1)
                            count += 1
                            if not replace_all:
                                break
                if count > 0 and not replace_all:
                    break

        # Also check tables
        if count == 0 or replace_all:
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.paragraphs:
                            if old_text in para.text:
                                for run in para.runs:
                                    if old_text in run.text:
                                        run.text = run.text.replace(old_text, new_text, 0 if replace_all else 1)
                                        count += 1
                                        if not replace_all:
                                            break

        if count == 0:
            return f"Error: '{old_text}' not found in {p.name}"

        doc.save(str(p))
        return f"Replaced {count} occurrence(s) in {p.name}"
    except Exception as e:
        return f"Error editing DOCX: {e}"


# ── PowerPoint (PPTX) ───────────────────────────────────────────────────────


@mcp.tool()
def read_pptx(path: str) -> str:
    """Extract text from a PowerPoint presentation.

    Args:
        path: PPTX file path (absolute or relative to workspace)
    """
    from pptx import Presentation

    p = _safe_path(path)
    if not p.exists():
        return f"Error: File not found: {p}"

    try:
        prs = Presentation(str(p))
        output = [f"# {p.name} ({len(prs.slides)} slides)\n"]

        for i, slide in enumerate(prs.slides, 1):
            output.append(f"--- Slide {i} ---")
            if slide.shapes.title:
                output.append(f"Title: {slide.shapes.title.text}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            output.append(text)
            output.append("")

        return "\n".join(output)
    except Exception as e:
        return f"Error reading PPTX: {e}"


@mcp.tool()
def create_pptx(path: str, slides: str, title: str = "") -> str:
    """Create a PowerPoint presentation.

    Args:
        path: Output PPTX path (absolute or relative to workspace)
        slides: JSON string — list of slide objects: [{"title": "...", "content": "..."}, ...]
        title: Optional presentation title (added as first slide)
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    p = _safe_path(path)
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        parsed = json.loads(slides)
        prs = Presentation()

        # Title slide
        if title:
            slide_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = ""

        # Content slides
        for item in parsed:
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = item.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = item.get("content", "")

            # Handle bullet points
            content = item.get("content", "")
            if "\n" in content:
                tf.clear()
                for line in content.split("\n"):
                    p_elem = tf.add_paragraph()
                    p_elem.text = line.strip().lstrip("- •")
                    p_elem.font.size = Pt(18)

        prs.save(str(p))
        return f"Created PPTX: {p} ({len(prs.slides)} slides)"
    except json.JSONDecodeError:
        return "Error: 'slides' must be valid JSON — list of {title, content} objects"
    except Exception as e:
        return f"Error creating PPTX: {e}"


if __name__ == "__main__":
    mcp.run(transport="stdio")
